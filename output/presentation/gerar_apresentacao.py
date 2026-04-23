"""Gera apresentacao PowerPoint do BioSPARQL-NL para publico leigo (25 slides, 15 min)."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# Paleta
AZUL   = RGBColor(0x1F, 0x38, 0x64)   # titulos
VERDE  = RGBColor(0x2E, 0x7D, 0x32)   # destaques
CINZA  = RGBColor(0x33, 0x33, 0x33)   # texto
BRANCO = RGBColor(0xFF, 0xFF, 0xFF)
FUNDO  = RGBColor(0xF5, 0xF7, 0xFA)   # fundo leve
AMARELO= RGBColor(0xFF, 0xC1, 0x07)   # aviso/destaque
LARANJA= RGBColor(0xE6, 0x51, 0x00)
VERMELHO=RGBColor(0xC6, 0x28, 0x28)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]  # completamente em branco


def slide_novo():
    return prs.slides.add_slide(BLANK)


def retangulo(slide, x, y, w, h, fill_rgb, radius=False):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    return shape


def texto(slide, txt, x, y, w, h, size=24, bold=False, color=CINZA,
          align=PP_ALIGN.LEFT, wrap=True, italic=False):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tb.word_wrap = wrap
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = txt
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return tb


def titulo_slide(slide, titulo, subtitulo=None, fundo_color=AZUL):
    retangulo(slide, 0, 0, 13.33, 7.5, BRANCO)
    retangulo(slide, 0, 0, 13.33, 1.6, fundo_color)
    texto(slide, titulo, 0.4, 0.25, 12.5, 1.2,
          size=36, bold=True, color=BRANCO, align=PP_ALIGN.LEFT)
    if subtitulo:
        texto(slide, subtitulo, 0.4, 1.7, 12.5, 0.6,
              size=20, color=AZUL, bold=False, align=PP_ALIGN.LEFT)


def bullets(slide, items, x=0.5, y=2.0, w=12.3, h=4.5, size=22, color=CINZA, spacing=0.55):
    for i, item in enumerate(items):
        prefixo = "• " if not item.startswith("    ") else "   – "
        t = prefixo + item.lstrip()
        texto(slide, t, x, y + i * spacing, w, spacing + 0.1,
              size=size, color=color)


def caixa_destaque(slide, txt, x, y, w, h, bg=VERDE, fg=BRANCO, size=20, bold=True):
    retangulo(slide, x, y, w, h, bg)
    texto(slide, txt, x + 0.15, y + 0.08, w - 0.3, h - 0.16,
          size=size, bold=bold, color=fg, align=PP_ALIGN.CENTER, wrap=True)


def tabela_simples(slide, cabecalho, linhas, x, y, col_widths, row_h=0.45,
                   header_bg=AZUL, alt_bg=RGBColor(0xE8, 0xEF, 0xF8)):
    # cabeçalho
    cx = x
    for i, (cel, w) in enumerate(zip(cabecalho, col_widths)):
        retangulo(slide, cx, y, w, row_h, header_bg)
        texto(slide, cel, cx + 0.05, y + 0.04, w - 0.1, row_h - 0.08,
              size=14, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
        cx += w
    # linhas
    for r, linha in enumerate(linhas):
        bg = alt_bg if r % 2 == 0 else BRANCO
        cx = x
        for i, (cel, w) in enumerate(zip(linha, col_widths)):
            retangulo(slide, cx, y + (r + 1) * row_h, w, row_h, bg)
            cor = VERDE if cel.startswith("+") else (VERMELHO if cel.startswith("-") else CINZA)
            texto(slide, cel, cx + 0.05, y + (r + 1) * row_h + 0.04,
                  w - 0.1, row_h - 0.08, size=13, color=cor, align=PP_ALIGN.CENTER)
            cx += w


def barra_horizontal(slide, label, valor, max_val, x, y, barra_w, barra_h, cor, size=14):
    pct = valor / max_val
    retangulo(slide, x + 2.2, y, barra_w, barra_h, RGBColor(0xE0, 0xE0, 0xE0))
    retangulo(slide, x + 2.2, y, barra_w * pct, barra_h, cor)
    texto(slide, label, x, y + 0.02, 2.1, barra_h, size=size, color=CINZA)
    texto(slide, f"{valor:.0f}%", x + 2.2 + barra_w * pct + 0.05, y + 0.02,
          0.8, barra_h, size=size, bold=True, color=cor)


# ============================================================
# SLIDE 1 — CAPA
# ============================================================
sl = slide_novo()
retangulo(sl, 0, 0, 13.33, 7.5, AZUL)
retangulo(sl, 0, 4.8, 13.33, 2.7, RGBColor(0x15, 0x26, 0x48))
texto(sl, "BioSPARQL-NL", 0.6, 1.0, 12.0, 1.6,
      size=54, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
texto(sl, "Como fazer perguntas a bancos de dados de doenças\nem linguagem natural", 0.6, 2.7, 12.0, 1.3,
      size=26, color=RGBColor(0xAA, 0xC4, 0xFF), align=PP_ALIGN.CENTER)
texto(sl, "Cassiano Ricardo Neubauer Moralles", 0.6, 4.95, 12.0, 0.5,
      size=18, color=BRANCO, align=PP_ALIGN.CENTER)
texto(sl, "Orientador: Prof. Sandro J. Rigo  |  PPG Computação Aplicada — UNISINOS  |  Abril 2026",
      0.6, 5.5, 12.0, 0.4, size=14, color=RGBColor(0xAA, 0xC4, 0xFF), align=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 2 — O PROBLEMA
# ============================================================
sl = slide_novo()
titulo_slide(sl, "O desafio da pesquisa biomédica")
texto(sl, "Imagine que você é médico e quer saber:", 0.5, 1.75, 12.3, 0.5,
      size=20, color=CINZA, italic=True)
caixa_destaque(sl, '"Quais doenças causam perda de memória e afetam crianças?"',
               1.0, 2.35, 11.3, 0.85, bg=RGBColor(0xE3, 0xF2, 0xFD),
               fg=AZUL, size=22, bold=False)
texto(sl, "Os dados existem — mas estão em bancos estruturados que exigem\nlinguagem técnica para serem consultados.", 0.5, 3.35, 12.3, 0.8,
      size=19, color=CINZA)
bullets(sl, [
    "Disease Ontology (DOID): banco com 12.000+ doenças catalogadas",
    "Human Phenotype Ontology (HPO): 16.000+ sintomas e fenótipos",
    "HPOA: 320.000+ associações doença-sintoma",
], size=19, y=4.25, spacing=0.6)

# ============================================================
# SLIDE 3 — O QUE SÃO ONTOLOGIAS
# ============================================================
sl = slide_novo()
titulo_slide(sl, "O que são essas bases de dados?")
texto(sl, "Pense numa enciclopédia médica ultra-estruturada:", 0.5, 1.75, 12.3, 0.5,
      size=20, color=CINZA)
# 3 caixas lado a lado
for i, (nome, desc, cor) in enumerate([
    ("DOID\n(Doenças)", "12.000+ doenças\norganizadas por família\n(ex: diabetes → doença metabólica)", AZUL),
    ("HPO\n(Sintomas)", "16.000+ fenótipos\n(ex: perda de memória,\nconvulsões, icterícia)", VERDE),
    ("HPOA\n(Associações)", "Quais doenças causam\nquais sintomas —\n320.000+ conexões", LARANJA),
]):
    caixa_destaque(sl, f"{nome}\n\n{desc}", 0.4 + i * 4.3, 2.4, 4.1, 3.0,
                   bg=cor, fg=BRANCO, size=17)
texto(sl, "Problema: consultar esses dados exige aprender uma linguagem técnica — SPARQL",
      0.5, 5.6, 12.3, 0.5, size=18, color=VERMELHO, bold=True, align=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 4 — A BARREIRA: SPARQL
# ============================================================
sl = slide_novo()
titulo_slide(sl, "A barreira técnica: SPARQL")
texto(sl, "Para consultar essas bases, um especialista precisaria escrever:", 0.5, 1.75, 12.3, 0.45, size=19, color=CINZA)
caixa_destaque(sl,
    "SELECT ?doenca ?label WHERE {\n"
    "  GRAPH <urn:hpoa> { ?a hpoa:has_phenotype obo:HP_0002354 }\n"
    "  GRAPH <urn:doid> { ?d rdfs:label ?label }\n"
    "  BIND(REPLACE(?xref,\"MIM:\",\"OMIM:\") AS ?omim)\n"
    "}",
    0.5, 2.3, 12.3, 2.0, bg=RGBColor(0x37, 0x47, 0x4F), fg=RGBColor(0xA5, 0xD6, 0xA7),
    size=16, bold=False)
texto(sl, "vs. o que o usuário quer digitar:", 0.5, 4.45, 12.3, 0.4, size=19, color=CINZA)
caixa_destaque(sl, '"Quais doenças causam perda de memória?"',
               0.5, 4.9, 12.3, 0.75, bg=RGBColor(0xE8, 0xF5, 0xE9), fg=VERDE, size=22, bold=False)
texto(sl, "O BioSPARQL-NL faz essa tradução automaticamente.",
      0.5, 5.8, 12.3, 0.45, size=20, bold=True, color=AZUL, align=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 5 — NOSSA PROPOSTA
# ============================================================
sl = slide_novo()
titulo_slide(sl, "Nossa proposta: BioSPARQL-NL")
# Fluxo visual simples
for i, (label, cor) in enumerate([
    ("Pergunta\nem português", RGBColor(0x42, 0xA5, 0xF5)),
    ("Pipeline\nBioSPARQL-NL", AZUL),
    ("Consulta\nSPARQL válida", VERDE),
    ("Resposta\ndo banco", RGBColor(0x66, 0xBB, 0x6A)),
]):
    caixa_destaque(sl, label, 0.5 + i * 3.2, 2.2, 2.8, 1.4, bg=cor, size=20)
    if i < 3:
        texto(sl, "→", 3.3 + i * 3.2, 2.7, 0.4, 0.6, size=32, bold=True, color=AZUL, align=PP_ALIGN.CENTER)
bullets(sl, [
    "Roda localmente no seu computador — seus dados nunca saem da sua máquina",
    "Não exige que o usuário saiba SPARQL",
    "Valida automaticamente se a consulta gerada é correta",
    "Corrige erros e tenta novamente se necessário",
], y=3.8, size=19, spacing=0.6)

# ============================================================
# SLIDE 6 — ARQUITETURA
# ============================================================
sl = slide_novo()
titulo_slide(sl, "Como funciona: arquitetura do pipeline")
img_path = os.path.join(os.path.dirname(__file__), "..", "relatorio", "arquitetura.png")
if os.path.exists(img_path):
    sl.shapes.add_picture(img_path, Inches(0.4), Inches(1.7), Inches(6.5), Inches(5.2))
bullets(sl, [
    "① Extrai entidades da pergunta (NER)",
    "② Busca exemplos similares (few-shot)",
    "③ Monta prompt com contexto",
    "④ LLM gera a consulta SPARQL",
    "⑤ Valida contra o esquema",
    "⑥ Executa no banco de dados",
], x=7.1, y=1.8, w=5.8, size=18, spacing=0.72)

# ============================================================
# SLIDE 7 — NER
# ============================================================
sl = slide_novo()
titulo_slide(sl, "Etapa 1: Entendendo a pergunta")
texto(sl, "O sistema usa Inteligência Artificial (scispaCy) para identificar\nentidades biomédicas na pergunta:", 0.5, 1.75, 12.3, 0.7, size=19, color=CINZA)
caixa_destaque(sl, '"Quais fenótipos estão associados à síndrome de Marfan?"',
               0.5, 2.6, 12.3, 0.75, bg=RGBColor(0xE3, 0xF2, 0xFD), fg=AZUL, size=21, bold=False)
texto(sl, "↓", 6.3, 3.45, 1.0, 0.5, size=28, bold=True, color=AZUL, align=PP_ALIGN.CENTER)
caixa_destaque(sl, "Reconhecido: \"síndrome de Marfan\"\n→ URI: obo:DOID_0060618  (no banco de doenças)",
               1.5, 4.05, 10.3, 1.0, bg=RGBColor(0xE8, 0xF5, 0xE9), fg=VERDE, size=19, bold=False)
texto(sl, "Isso evita que a IA precise adivinhar como a doença está cadastrada no sistema.",
      0.5, 5.25, 12.3, 0.5, size=17, color=CINZA, italic=True)

# ============================================================
# SLIDE 8 — FEW-SHOT / RAG
# ============================================================
sl = slide_novo()
titulo_slide(sl, "Etapa 2: Aprendendo com exemplos")
texto(sl, "O sistema mantém uma biblioteca de 52 perguntas + consultas já validadas.", 0.5, 1.75, 12.3, 0.5, size=19, color=CINZA)
texto(sl, "Quando chega uma nova pergunta, busca as 3 mais parecidas:", 0.5, 2.3, 12.3, 0.45, size=19, color=CINZA)
for i, (q, sim) in enumerate([
    ('"Sintomas do Alzheimer?"', "92% similar"),
    ('"Fenótipos da doença de Huntington?"', "87% similar"),
    ('"Manifestações clínicas da esclerose múltipla?"', "81% similar"),
]):
    caixa_destaque(sl, f"{q}  [{sim}]", 0.5, 2.95 + i * 0.82, 12.3, 0.7,
                   bg=RGBColor(0xE8, 0xEF, 0xF8), fg=AZUL, size=17, bold=False)
texto(sl, "Esses exemplos são incluídos no prompt da IA — ela aprende o padrão\nantes de gerar a nova consulta. Isso reduz alucinações.", 0.5, 5.55, 12.3, 0.7, size=18, color=CINZA)

# ============================================================
# SLIDE 9 — GERAÇÃO
# ============================================================
sl = slide_novo()
titulo_slide(sl, "Etapa 3: A IA gera a consulta")
texto(sl, "O prompt enviado à IA contém:", 0.5, 1.75, 12.3, 0.45, size=19, color=CINZA)
bullets(sl, [
    "Instruções sobre a estrutura dos bancos de dados",
    "Esquema: quais propriedades e classes existem",
    "Exemplos de perguntas similares já resolvidas",
    "A pergunta do usuário",
], y=2.35, size=18, spacing=0.55)
texto(sl, "Modelos testados (rodam localmente no seu PC):", 0.5, 4.65, 12.3, 0.45, size=19, color=CINZA, bold=True)
for i, (nome, params, cor) in enumerate([
    ("Gemma 3", "4B params", RGBColor(0x42, 0xA5, 0xF5)),
    ("Nemotron Nano", "4B params", VERDE),
    ("Qwen 3.5", "9B params", LARANJA),
    ("Claude Sonnet*", "comercial", VERMELHO),
]):
    caixa_destaque(sl, f"{nome}\n{params}", 0.3 + i * 3.2, 5.2, 2.9, 1.6, bg=cor, size=17)
texto(sl, "* Claude Sonnet rodou via API como referência de teto de desempenho",
      0.3, 6.9, 12.7, 0.4, size=13, color=CINZA, italic=True)

# ============================================================
# SLIDE 10 — VALIDAÇÃO
# ============================================================
sl = slide_novo()
titulo_slide(sl, "Etapa 4: Validação e autocorreção")
texto(sl, "Problema: LLMs às vezes inventam propriedades que não existem no banco.", 0.5, 1.75, 12.3, 0.5, size=19, color=VERMELHO)
# Fluxo de validação
caixa_destaque(sl, "SPARQL gerado\npela IA", 0.3, 2.4, 2.8, 1.2, bg=AZUL, size=17)
texto(sl, "→", 3.1, 2.7, 0.6, 0.6, size=28, bold=True, color=CINZA, align=PP_ALIGN.CENTER)
caixa_destaque(sl, "Validação\npor esquema", 3.7, 2.4, 2.8, 1.2, bg=RGBColor(0x78, 0x90, 0x9C), size=17)
texto(sl, "→", 6.5, 2.7, 0.6, 0.6, size=28, bold=True, color=VERDE, align=PP_ALIGN.CENTER)
caixa_destaque(sl, "✓ Válida\n→ executa", 7.1, 2.4, 2.4, 1.2, bg=VERDE, size=17)
texto(sl, "↺", 6.5, 3.7, 0.6, 0.6, size=28, bold=True, color=VERMELHO, align=PP_ALIGN.CENTER)
caixa_destaque(sl, "✗ Inválida\n→ corrige (até 2x)", 7.1, 3.65, 2.4, 1.2, bg=VERMELHO, size=16)
texto(sl, "A validação usa introspecção: o sistema extrai automaticamente do próprio banco\nquais propriedades e classes existem — sem precisar de configuração manual.", 0.5, 5.1, 12.3, 0.8, size=18, color=CINZA)
caixa_destaque(sl, "Qwen 9B: de 0% → 50% de acerto só com esse mecanismo (+50pp)!",
               0.5, 6.0, 12.3, 0.75, bg=RGBColor(0xE8, 0xF5, 0xE9), fg=VERDE, size=19)

# ============================================================
# SLIDE 11 — EXEMPLO REAL
# ============================================================
sl = slide_novo()
titulo_slide(sl, "Exemplo real de funcionamento")
caixa_destaque(sl, "Pergunta: \"What phenotypes are associated with Marfan syndrome?\"",
               0.3, 1.75, 12.7, 0.7, bg=RGBColor(0xE3, 0xF2, 0xFD), fg=AZUL, size=19, bold=False)
texto(sl, "↓  SPARQL gerado e executado:", 0.5, 2.55, 12.3, 0.4, size=17, color=CINZA)
caixa_destaque(sl,
    "SELECT DISTINCT ?phenoLabel WHERE {\n"
    "  GRAPH <urn:hpoa> { ?ann hpoa:source_id \"OMIM:154700\" ;\n"
    "                         hpoa:has_phenotype ?pheno }\n"
    "  GRAPH <urn:hpo> { ?pheno rdfs:label ?phenoLabel }\n"
    "} LIMIT 10",
    0.3, 3.0, 12.7, 1.8, bg=RGBColor(0x37, 0x47, 0x4F), fg=RGBColor(0xA5, 0xD6, 0xA7),
    size=15, bold=False)
caixa_destaque(sl,
    "Resultados: Arachnodactyly, Ectopia lentis, Aortic root aneurysm,\nScoliosis, Tall stature, Pectus excavatum ...",
    0.3, 4.9, 12.7, 0.9, bg=RGBColor(0xE8, 0xF5, 0xE9), fg=VERDE, size=18, bold=False)
texto(sl, "O usuário perguntou em inglês simples — o sistema consultou 3 bancos de dados e retornou os fenótipos.",
      0.3, 5.95, 12.7, 0.5, size=17, color=CINZA, italic=True)

# ============================================================
# SLIDE 12 — GOLD STANDARD / AVALIAÇÃO
# ============================================================
sl = slide_novo()
titulo_slide(sl, "Como avaliamos o sistema")
texto(sl, "Criamos um gold standard: 52 perguntas com respostas esperadas verificadas por 2 anotadores.", 0.5, 1.75, 12.3, 0.55, size=19, color=CINZA)
for i, (nivel, n, desc, cor) in enumerate([
    ("Fácil",  "13 perguntas", "1 grafo, 1 entidade\nEx: sintomas de diabetes", RGBColor(0x66, 0xBB, 0x6A)),
    ("Médio",  "18 perguntas", "2 grafos, múltiplas entidades\nEx: doenças com 5+ sintomas", LARANJA),
    ("Difícil","21 perguntas", "3 grafos, joins complexos,\nfiltros MIM/OMIM específicos", VERMELHO),
]):
    caixa_destaque(sl, f"{nivel}\n{n}\n\n{desc}", 0.3 + i * 4.35, 2.5, 4.05, 3.2,
                   bg=cor, fg=BRANCO, size=16)
texto(sl, "Métrica principal: 'correção' — a consulta retornou os resultados esperados?",
      0.3, 5.9, 12.7, 0.5, size=18, bold=True, color=AZUL, align=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 13 — MODELOS
# ============================================================
sl = slide_novo()
titulo_slide(sl, "Os 4 modelos de IA testados")
tabela_simples(sl,
    ["Modelo", "Parâmetros", "Onde roda", "Custo"],
    [
        ["Gemma 3 4B",      "4 bilhões",  "Seu PC (GPU 8GB)", "Gratuito"],
        ["Nemotron Nano 4B","4 bilhões",  "Seu PC (GPU 8GB)", "Gratuito"],
        ["Qwen 3.5 9B",     "9 bilhões",  "Seu PC (CPU+GPU)", "Gratuito"],
        ["Claude Sonnet",   "não divulgado", "Nuvem (API)",   "Pago"],
    ],
    x=0.5, y=1.85, col_widths=[3.2, 2.5, 3.3, 2.8],
    row_h=0.6
)
texto(sl, "Os três primeiros rodam localmente — nenhum dado sai do seu computador.\nClaude Sonnet serve como referência de teto (o melhor possível).", 0.5, 5.45, 12.3, 0.8, size=18, color=CINZA)
caixa_destaque(sl, "Descoberta: um modelo de 4B parâmetros superou o de 9B — tamanho não é tudo!",
               0.5, 6.35, 12.3, 0.75, bg=RGBColor(0xE8, 0xF5, 0xE9), fg=VERDE, size=18)

# ============================================================
# SLIDE 14 — RESULTADOS GERAIS 30Q
# ============================================================
sl = slide_novo()
titulo_slide(sl, "Resultados: 30 questões primárias")
texto(sl, "Com validação ativada (melhor cenário para cada modelo):", 0.5, 1.75, 12.3, 0.45, size=19, color=CINZA)
# barras horizontais
cores_bar = [VERDE, RGBColor(0x42, 0xA5, 0xF5), LARANJA]
for i, (label, val, cor) in enumerate([
    ("Nemotron Nano 4B", 80, VERDE),
    ("Qwen 3.5 9B",      50, LARANJA),
    ("Gemma 3 4B",       47, RGBColor(0x42, 0xA5, 0xF5)),
]):
    barra_horizontal(sl, label, val, 100, 0.5, 2.4 + i * 0.95, 9.5, 0.72, cor, size=18)
caixa_destaque(sl, "Conclusão surpreendente:\nNemotron 4B (80%) > Qwen 9B (50%)\nArquitetura importa mais que tamanho!",
               0.5, 5.35, 12.3, 1.5, bg=AZUL, size=20)

# ============================================================
# SLIDE 15 — IMPACTO DA VALIDAÇÃO
# ============================================================
sl = slide_novo()
titulo_slide(sl, "A validação faz diferença?")
texto(sl, "Comparação: com e sem o mecanismo de autocorreção por esquema", 0.5, 1.75, 12.3, 0.45, size=19, color=CINZA)
tabela_simples(sl,
    ["Modelo", "Sem validação", "Com validação", "Ganho"],
    [
        ["Gemma 3 4B",      "46,7%",  "46,7%",  "0,0pp"],
        ["Nemotron Nano 4B","86,7%",  "80,0%",  "-6,7pp"],
        ["Qwen 3.5 9B",     " 0,0%",  "50,0%",  "+50,0pp"],
    ],
    x=0.5, y=2.3, col_widths=[3.5, 2.8, 2.8, 2.2], row_h=0.65
)
bullets(sl, [
    "Qwen: validação essencial — de 0% a 50% de acerto (sem ela, inútil!)",
    "Gemma: sem diferença — já gera consultas que passam na validação",
    "Nemotron: leve redução — valida corretamente mas rejeita alguns queries válidos",
], y=5.1, size=18, spacing=0.58)

# ============================================================
# SLIDE 16 — POR DIFICULDADE
# ============================================================
sl = slide_novo()
titulo_slide(sl, "Resultado por nível de dificuldade")
tabela_simples(sl,
    ["Modelo", "Fácil (10q)", "Médio (10q)", "Difícil (10q)", "Geral"],
    [
        ["Gemma 3 4B",      "60%", "40%", "40%", "47%"],
        ["Nemotron Nano 4B","80%", "80%", "80%", "80%"],
        ["Qwen 3.5 9B",     "50%", "30%", "70%", "50%"],
    ],
    x=0.4, y=1.85, col_widths=[3.2, 2.3, 2.3, 2.3, 2.0], row_h=0.65
)
texto(sl, "Observações:", 0.5, 4.65, 12.3, 0.4, size=18, bold=True, color=AZUL)
bullets(sl, [
    "Nemotron: desempenho uniforme (80% em todos os níveis) — mais confiável",
    "Qwen: melhor em difícil (70%) — bom em queries complexas com OPTIONAL/FILTER",
    "Gemma: mais rápido (4s/query) mas menos preciso",
], y=5.15, size=17, spacing=0.58)

# ============================================================
# SLIDE 17 — LOCAL VS COMERCIAL
# ============================================================
sl = slide_novo()
titulo_slide(sl, "Local vs. Comercial: 52 questões")
texto(sl, "Avaliamos o Claude Sonnet (modelo comercial de referência) no gold standard completo,\njunto com o melhor modelo local (Gemma 3 4B) para comparação:", 0.5, 1.75, 12.3, 0.7, size=18, color=CINZA)
tabela_simples(sl,
    ["Modelo", "Geral", "Fácil", "Médio", "Difícil", "Sintaxe"],
    [
        ["Gemma 3 4B (local)", "46,2%", "39%", "56%", "43%", "75%"],
        ["Claude Sonnet (nuvem)", "82,7%", "100%", "89%", "67%", "100%"],
    ],
    x=0.3, y=2.6, col_widths=[3.5, 1.8, 1.7, 1.7, 1.7, 1.6], row_h=0.65
)
texto(sl, "Claude: 43/52 corretas  ·  Gemma: 24/52 corretas", 0.3, 4.2, 12.7, 0.45,
      size=17, color=CINZA, align=PP_ALIGN.CENTER)
caixa_destaque(sl, "Gap: 36,5 pontos percentuais", 2.0, 4.75, 9.3, 0.8,
               bg=AZUL, fg=BRANCO, size=22)
texto(sl, "O Claude Sonnet tem 100% de sintaxe válida — nunca gera SPARQL quebrado.",
      0.3, 5.7, 12.7, 0.45, size=17, color=CINZA, italic=True)

# ============================================================
# SLIDE 18 — GAP TRADE-OFF
# ============================================================
sl = slide_novo()
titulo_slide(sl, "O trade-off: privacidade vs precisão")
# Visualização do gap
retangulo(sl, 0.5, 1.8, 12.3, 1.0, RGBColor(0xE3, 0xF2, 0xFD))
retangulo(sl, 0.5, 1.8, 12.3 * 0.462, 1.0, RGBColor(0x42, 0xA5, 0xF5))
texto(sl, "Gemma local  46,2%", 0.7, 2.0, 5.0, 0.6, size=18, bold=True, color=BRANCO)
retangulo(sl, 0.5, 3.1, 12.3, 1.0, RGBColor(0xE8, 0xF5, 0xE9))
retangulo(sl, 0.5, 3.1, 12.3 * 0.827, 1.0, VERDE)
texto(sl, "Claude Sonnet  82,7%", 0.7, 3.3, 6.0, 0.6, size=18, bold=True, color=BRANCO)
bullets(sl, [
    "Gemma local: dados ficam no seu PC, privacidade total — mas 36pp a menos",
    "Claude Sonnet: máxima precisão — mas seus dados vão para a nuvem (API paga)",
    "Para dados médicos sensíveis, o modelo local pode ser a escolha certa",
    "Um PC com GPU de 8GB já é suficiente para rodar o Gemma 4B",
], y=4.3, size=18, spacing=0.6)

# ============================================================
# SLIDE 19 — ABLAÇÃO
# ============================================================
sl = slide_novo()
titulo_slide(sl, "O que cada componente contribui?")
texto(sl, "Estudo de ablação: removemos um componente por vez e medimos o impacto\n(Gemma 3 4B, 30 questões)", 0.5, 1.75, 12.3, 0.65, size=18, color=CINZA)
tabela_simples(sl,
    ["Configuração", "Sintaxe", "Correção", "Δ vs completo"],
    [
        ["Pipeline completo",   "90%", "50,0%", "---"],
        ["Sem NER",             "90%", "63,3%", "+13,3pp"],
        ["Sem few-shot",        "90%", "46,7%", "-3,3pp"],
        ["Sem esquema",         "77%", "46,7%", "-3,3pp"],
        ["Sem validação",       "67%", "50,0%", "0,0pp"],
        ["Zero-shot (nada)",    "87%", "66,7%", "+16,7pp"],
    ],
    x=0.3, y=2.55, col_widths=[3.8, 2.0, 2.0, 3.2], row_h=0.52
)
texto(sl, "* NER (reconhecimento de entidades) piora o Gemma — injeta ruído para este modelo específico.",
      0.3, 6.6, 12.7, 0.45, size=13, color=CINZA, italic=True)

# ============================================================
# SLIDE 20 — INSIGHT: VALIDAÇÃO CRÍTICA
# ============================================================
sl = slide_novo()
titulo_slide(sl, "Achado 1: Validação é o componente crítico")
texto(sl, "O loop de autocorreção por esquema é essencial — veja o impacto sem ele:", 0.5, 1.75, 12.3, 0.5, size=19, color=CINZA)
for i, (modelo, com_v, sem_v, delta, cor) in enumerate([
    ("Gemma 3 4B",   "46,2%", "25,0%", "−21,2pp", RGBColor(0x42, 0xA5, 0xF5)),
    ("Claude Sonnet","82,7%", "46,2%", "−36,5pp", VERDE),
]):
    retangulo(sl, 0.4, 2.45 + i * 2.2, 12.5, 1.8, RGBColor(0xF5, 0xF7, 0xFA))
    texto(sl, modelo, 0.6, 2.5 + i * 2.2, 3.5, 0.5, size=20, bold=True, color=AZUL)
    caixa_destaque(sl, f"Com validação\n{com_v}", 0.5, 2.95 + i * 2.2, 3.8, 1.1, bg=cor, size=19)
    texto(sl, "→", 4.4, 3.25 + i * 2.2, 0.8, 0.6, size=30, bold=True, color=VERMELHO, align=PP_ALIGN.CENTER)
    caixa_destaque(sl, f"Sem validação\n{sem_v}", 5.2, 2.95 + i * 2.2, 3.8, 1.1, bg=VERMELHO, size=19)
    caixa_destaque(sl, delta, 9.2, 2.95 + i * 2.2, 3.2, 1.1, bg=RGBColor(0x37, 0x47, 0x4F), size=22)

# ============================================================
# SLIDE 21 — INSIGHT: ARQUITETURA > TAMANHO
# ============================================================
sl = slide_novo()
titulo_slide(sl, "Achado 2: Arquitetura importa mais que tamanho")
texto(sl, "Comparando modelos de tamanhos diferentes:", 0.5, 1.75, 12.3, 0.45, size=19, color=CINZA)
for i, (nome, params, acerto, cor) in enumerate([
    ("Nemotron Nano", "4 bilhões", 80, VERDE),
    ("Qwen 3.5",      "9 bilhões", 50, LARANJA),
    ("Gemma 3",       "4 bilhões", 47, RGBColor(0x42, 0xA5, 0xF5)),
]):
    barra_horizontal(sl, f"{nome} ({params})", acerto, 100, 0.5, 2.4 + i * 1.05, 9.3, 0.78, cor, size=18)
texto(sl, "Nemotron Nano tem 4B parâmetros (MENOS que o Qwen 9B) mas acerta 80% vs 50%",
      0.5, 5.7, 12.3, 0.5, size=19, color=CINZA)
caixa_destaque(sl,
    "Conclusão: a arquitetura e o treinamento do modelo são mais determinantes\nque o número de parâmetros para gerar SPARQL biomédico correto.",
    0.5, 6.3, 12.3, 0.9, bg=AZUL, size=18)

# ============================================================
# SLIDE 22 — INTERFACE XAI
# ============================================================
sl = slide_novo()
titulo_slide(sl, "Interface explicável (XAI)")
bullets(sl, [
    "Interface web desenvolvida em React + FastAPI",
    "O usuário digita a pergunta em linguagem natural",
    "Cada etapa do pipeline é mostrada em tempo real:",
    "    → Entidades reconhecidas (NER)",
    "    → Exemplos recuperados (few-shot)",
    "    → Consulta SPARQL gerada",
    "    → Resultado da validação e tentativas",
    "    → Resultado final do banco de dados",
    "Alinhada com princípios de IA Explicável (XAI): nada é 'caixa preta'",
], y=1.85, size=17, spacing=0.5)

# ============================================================
# SLIDE 23 — LIMITAÇÕES
# ============================================================
sl = slide_novo()
titulo_slide(sl, "O que ainda não funciona bem")
for i, (titulo, desc, cor) in enumerate([
    ("Hardware como gargalo",
     "Modelos ≤4B rodam bem em GPU 8GB. O Qwen 9B é 15x mais lento (63s por query). Modelos maiores (>20B) não cabem em hardware comum.",
     AZUL),
    ("Convenções idiossincráticas",
     "9 perguntas (17%) falham mesmo com o Claude Sonnet. Envolvem incompatibilidades técnicas entre bases (ex: DOID usa 'MIM:' mas HPOA usa 'OMIM:').",
     LARANJA),
    ("Variabilidade estocástica",
     "O Gemma 3 4B variou ±3 acertos entre execuções idênticas (temperature=0.1). Resultados não são 100% determinísticos.",
     VERMELHO),
]):
    retangulo(sl, 0.3, 1.85 + i * 1.8, 12.7, 1.6, RGBColor(0xF5, 0xF7, 0xFA))
    caixa_destaque(sl, titulo, 0.3, 1.85 + i * 1.8, 3.2, 1.6, bg=cor, size=16)
    texto(sl, desc, 3.65, 2.0 + i * 1.8, 9.2, 1.3, size=16, color=CINZA)

# ============================================================
# SLIDE 24 — TABELA EXPANDIDA: 7 MODELOS
# ============================================================
sl = slide_novo()
titulo_slide(sl, "Comparação: 7 modelos (52 questões)")
modelos = [
    ("BioGPT 1.5B",      " 0,0%", " 0,0%", "GPT-2 style — sem instruction-following"),
    ("BioMedLM 2.7B",    " 0,0%", " 0,0%", "GPT-2 style — sem instruction-following"),
    ("Gemma 3 4B",       "46,2%", "75,0%", "Modelo geral (baseline)"),
    ("OpenBioLLM-8B",    "57,7%", "61,5%", "Biomédico — instruct fraco"),
    ("Llama3-Med42-8B",  "67,3%", "90,4%", "Biomédico — instruct forte"),
    ("Qwen3-VL-8B",      "78,8%", "98,1%", "Geral multimodal — melhor local"),
    ("Claude Sonnet",    "82,7%","100,0%", "Upper bound comercial"),
]
cores_linha = [CINZA, CINZA, AZUL, VERDE, VERDE, LARANJA, VERMELHO]
retangulo(sl, 0.3, 1.3, 12.7, 0.45, AZUL)
for col, cab in enumerate(["Modelo", "Correção", "Sintaxe", "Observação"]):
    xs = [0.35, 5.1, 6.7, 8.1][col]
    ws = [4.6, 1.4, 1.3, 4.8][col]
    texto(sl, cab, xs, 1.32, ws, 0.42, size=13, bold=True, color=BRANCO)
for i, (nome, corr, sint, obs) in enumerate(modelos):
    y = 1.78 + i * 0.72
    bg = RGBColor(0xF0, 0xF4, 0xFF) if i % 2 == 0 else RGBColor(0xFF, 0xFF, 0xFF)
    retangulo(sl, 0.3, y, 12.7, 0.68, bg)
    cor = cores_linha[i]
    texto(sl, nome, 0.35, y + 0.05, 4.6, 0.55, size=13, bold=True, color=cor)
    texto(sl, corr, 5.1,  y + 0.05, 1.4, 0.55, size=13, bold=True, color=cor, align=PP_ALIGN.CENTER)
    texto(sl, sint, 6.7,  y + 0.05, 1.3, 0.55, size=13, color=CINZA, align=PP_ALIGN.CENTER)
    texto(sl, obs,  8.1,  y + 0.05, 4.8, 0.55, size=12, color=CINZA)

# ============================================================
# SLIDE 25 — QALD-9+: GENERALIZAÇÃO
# ============================================================
sl = slide_novo()
titulo_slide(sl, "Generalização: QALD-9+ (DBpedia)")
retangulo(sl, 0.4, 1.3, 12.5, 1.5, RGBColor(0xE3, 0xF2, 0xFD))
texto(sl, "Pergunta: o pipeline funciona em outros domínios além do biomédico?",
      0.6, 1.4, 12.0, 0.5, size=17, bold=True, color=AZUL)
texto(sl, "Testamos o Gemma 3 4B sobre 30 perguntas do benchmark público QALD-9+ (DBpedia) — sem adaptação de domínio.",
      0.6, 1.85, 12.0, 0.45, size=14, color=CINZA)
for i, (label, valor, cor, desc) in enumerate([
    ("Execução OK", "80%",   VERDE,   "SPARQL enviado ao DBpedia e aceito"),
    ("F1 médio",    "0,106", LARANJA, "Poucos resultados corretos"),
    ("Exact match", "3,3%",  VERMELHO,"Apenas 1/30 perfeito"),
]):
    x = 0.5 + i * 4.2
    retangulo(sl, x, 2.95, 3.8, 1.6, RGBColor(0xF5, 0xF5, 0xF5))
    texto(sl, valor, x + 0.15, 3.05, 3.5, 0.9, size=48, bold=True, color=cor, align=PP_ALIGN.CENTER)
    texto(sl, label, x + 0.15, 3.9,  3.5, 0.4, size=15, bold=True, color=AZUL, align=PP_ALIGN.CENTER)
    texto(sl, desc,  x + 0.15, 4.3,  3.5, 0.35, size=12, color=CINZA, align=PP_ALIGN.CENTER)
retangulo(sl, 0.4, 5.1, 12.5, 1.7, RGBColor(0xFF, 0xF3, 0xE0))
texto(sl, "Conclusão:", 0.6, 5.2, 3.0, 0.45, size=15, bold=True, color=LARANJA)
texto(sl, "O pipeline gera SPARQL válido em outros domínios (80% execução), mas precisa de exemplos e esquema\n"
          "específicos do domínio alvo para acertar semanticamente. A transferência sintática funciona; a semântica não.",
      0.6, 5.6, 12.0, 1.1, size=14, color=CINZA)

# ============================================================
# SLIDE 26 — CONCLUSÃO
# ============================================================
sl = slide_novo()
titulo_slide(sl, "5 achados principais")
for i, (letra, titulo, desc, cor) in enumerate([
    ("a", "Validação é essencial",
     "Gemma: 46% → 25% sem ela; Claude: 83% → 46% — o loop de autocorreção é o componente mais crítico",
     AZUL),
    ("b", "Arquitetura > tamanho",
     "Nemotron 4B (80%) supera Qwen 9B (50%) — treinamento importa mais que número de parâmetros",
     VERDE),
    ("c", "Gap quase fechado com modelo certo",
     "Qwen3-VL 78,8% vs Claude 82,7% (gap de 3,9pp) — modelos locais de próxima geração viabilizam privacidade + precisão",
     LARANJA),
    ("d", "Erros irredutíveis existem",
     "9/52 perguntas (17%) falham mesmo com o melhor modelo — limitações intrínsecas dos dados",
     VERMELHO),
    ("e", "Domínio biomédico ajuda, mas instruction-following é obrigatório",
     "Llama3-Med42 (67%) supera Gemma 4B (46%); BioGPT/BioMedLM GPT-2 style = 0%",
     RGBColor(0x6A, 0x1B, 0x9A)),
]):
    retangulo(sl, 0.3, 1.75 + i * 1.08, 0.65, 0.9, cor)
    texto(sl, letra, 0.3, 1.8 + i * 1.08, 0.65, 0.9, size=22, bold=True,
          color=BRANCO, align=PP_ALIGN.CENTER)
    texto(sl, titulo, 1.05, 1.78 + i * 1.08, 3.5, 0.45, size=15, bold=True, color=AZUL)
    texto(sl, desc, 1.05, 2.2 + i * 1.08, 12.0, 0.42, size=13, color=CINZA)

# ============================================================
# SLIDE 25 — OBRIGADO
# ============================================================
sl = slide_novo()
retangulo(sl, 0, 0, 13.33, 7.5, AZUL)
retangulo(sl, 0, 4.5, 13.33, 3.0, RGBColor(0x15, 0x26, 0x48))
texto(sl, "Obrigado!", 0.6, 1.1, 12.0, 1.3,
      size=60, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
texto(sl, "Perguntas?", 0.6, 2.5, 12.0, 0.9,
      size=34, color=RGBColor(0xAA, 0xC4, 0xFF), align=PP_ALIGN.CENTER)
texto(sl, "Cassiano Ricardo Neubauer Moralles",
      0.6, 4.7, 12.0, 0.5, size=18, color=BRANCO, align=PP_ALIGN.CENTER)
texto(sl, "cass78@gmail.com  |  PPG Computação Aplicada — UNISINOS",
      0.6, 5.25, 12.0, 0.45, size=15, color=RGBColor(0xAA, 0xC4, 0xFF), align=PP_ALIGN.CENTER)
texto(sl, "Código e dados: output/eval_*.json  |  output/ablation_*.json  |  output/relatorio/biosparql-nl.pdf",
      0.6, 5.75, 12.0, 0.4, size=12, color=RGBColor(0x78, 0x90, 0x9C), align=PP_ALIGN.CENTER)

# ============================================================
out = os.path.join(os.path.dirname(__file__), "biosparql-nl-apresentacao.pptx")
prs.save(out)
print(f"[OK] Apresentacao salva em: {out}")
print(f"     Total de slides: {len(prs.slides)}")
