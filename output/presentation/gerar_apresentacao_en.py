"""Generates BioSPARQL-NL PowerPoint presentation in English (27 slides, 15 min)."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# Color palette
AZUL   = RGBColor(0x1F, 0x38, 0x64)   # titles
VERDE  = RGBColor(0x2E, 0x7D, 0x32)   # highlights
CINZA  = RGBColor(0x33, 0x33, 0x33)   # text
BRANCO = RGBColor(0xFF, 0xFF, 0xFF)
FUNDO  = RGBColor(0xF5, 0xF7, 0xFA)   # light background
AMARELO= RGBColor(0xFF, 0xC1, 0x07)   # warning/highlight
LARANJA= RGBColor(0xE6, 0x51, 0x00)
VERMELHO=RGBColor(0xC6, 0x28, 0x28)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]  # completely blank


def slide_novo():
    return prs.slides.add_slide(BLANK)


def retangulo(slide, x, y, w, h, fill_rgb, radius=False):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    return shape


def texto(slide, txt, x, y, w, h, size=24, bold=False, color=CINZA,
          align=PP_ALIGN.LEFT, wrap=True, italic=False):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tb.word_wrap = wrap
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = txt
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return tb


def titulo_slide(slide, titulo, subtitulo=None, fundo_color=AZUL):
    retangulo(slide, 0, 0, 13.33, 7.5, BRANCO)
    retangulo(slide, 0, 0, 13.33, 1.6, fundo_color)
    texto(slide, titulo, 0.4, 0.25, 12.5, 1.2,
          size=36, bold=True, color=BRANCO, align=PP_ALIGN.LEFT)
    if subtitulo:
        texto(slide, subtitulo, 0.4, 1.7, 12.5, 0.6,
              size=20, color=AZUL, bold=False, align=PP_ALIGN.LEFT)


def bullets(slide, items, x=0.5, y=2.0, w=12.3, h=4.5, size=22, color=CINZA, spacing=0.55):
    for i, item in enumerate(items):
        prefixo = "• " if not item.startswith("    ") else "   – "
        t = prefixo + item.lstrip()
        texto(slide, t, x, y + i * spacing, w, spacing + 0.1,
              size=size, color=color)


def caixa_destaque(slide, txt, x, y, w, h, bg=VERDE, fg=BRANCO, size=20, bold=True):
    retangulo(slide, x, y, w, h, bg)
    texto(slide, txt, x + 0.15, y + 0.08, w - 0.3, h - 0.16,
          size=size, bold=bold, color=fg, align=PP_ALIGN.CENTER, wrap=True)


def tabela_simples(slide, cabecalho, linhas, x, y, col_widths, row_h=0.45,
                   header_bg=AZUL, alt_bg=RGBColor(0xE8, 0xEF, 0xF8)):
    # header
    cx = x
    for i, (cel, w) in enumerate(zip(cabecalho, col_widths)):
        retangulo(slide, cx, y, w, row_h, header_bg)
        texto(slide, cel, cx + 0.05, y + 0.04, w - 0.1, row_h - 0.08,
              size=14, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
        cx += w
    # rows
    for r, linha in enumerate(linhas):
        bg = alt_bg if r % 2 == 0 else BRANCO
        cx = x
        for i, (cel, w) in enumerate(zip(linha, col_widths)):
            retangulo(slide, cx, y + (r + 1) * row_h, w, row_h, bg)
            cor = VERDE if cel.startswith("+") else (VERMELHO if cel.startswith("-") else CINZA)
            texto(slide, cel, cx + 0.05, y + (r + 1) * row_h + 0.04,
                  w - 0.1, row_h - 0.08, size=13, color=cor, align=PP_ALIGN.CENTER)
            cx += w


def barra_horizontal(slide, label, valor, max_val, x, y, barra_w, barra_h, cor, size=14):
    pct = valor / max_val
    retangulo(slide, x + 2.2, y, barra_w, barra_h, RGBColor(0xE0, 0xE0, 0xE0))
    retangulo(slide, x + 2.2, y, barra_w * pct, barra_h, cor)
    texto(slide, label, x, y + 0.02, 2.1, barra_h, size=size, color=CINZA)
    texto(slide, f"{valor:.0f}%", x + 2.2 + barra_w * pct + 0.05, y + 0.02,
          0.8, barra_h, size=size, bold=True, color=cor)


# ============================================================
# SLIDE 1 — COVER
# ============================================================
sl = slide_novo()
retangulo(sl, 0, 0, 13.33, 7.5, AZUL)
retangulo(sl, 0, 4.8, 13.33, 2.7, RGBColor(0x15, 0x26, 0x48))
texto(sl, "BioSPARQL-NL", 0.6, 1.0, 12.0, 1.6,
      size=54, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
texto(sl, "How to query biomedical disease databases\nusing natural language", 0.6, 2.7, 12.0, 1.3,
      size=26, color=RGBColor(0xAA, 0xC4, 0xFF), align=PP_ALIGN.CENTER)
texto(sl, "Cassiano Ricardo Neubauer Moralles", 0.6, 4.95, 12.0, 0.5,
      size=18, color=BRANCO, align=PP_ALIGN.CENTER)
texto(sl, "Advisor: Prof. Sandro J. Rigo  |  Graduate Program in Applied Computing — UNISINOS  |  April 2026",
      0.6, 5.5, 12.0, 0.4, size=14, color=RGBColor(0xAA, 0xC4, 0xFF), align=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 2 — THE PROBLEM
# ============================================================
sl = slide_novo()
titulo_slide(sl, "The Biomedical Research Challenge")
texto(sl, "Imagine you are a physician and want to know:", 0.5, 1.75, 12.3, 0.5,
      size=20, color=CINZA, italic=True)
caixa_destaque(sl, '"Which diseases cause memory loss and affect children?"',
               1.0, 2.35, 11.3, 0.85, bg=RGBColor(0xE3, 0xF2, 0xFD),
               fg=AZUL, size=22, bold=False)
texto(sl, "The data exists — but is stored in structured databases that require\ntechnical language to query.", 0.5, 3.35, 12.3, 0.8,
      size=19, color=CINZA)
bullets(sl, [
    "Disease Ontology (DOID): database with 12,000+ catalogued diseases",
    "Human Phenotype Ontology (HPO): 16,000+ symptoms and phenotypes",
    "HPOA: 320,000+ disease-phenotype associations",
], size=19, y=4.25, spacing=0.6)

# ============================================================
# SLIDE 3 — WHAT ARE ONTOLOGIES
# ============================================================
sl = slide_novo()
titulo_slide(sl, "What Are These Databases?")
texto(sl, "Think of a highly structured medical encyclopedia:", 0.5, 1.75, 12.3, 0.5,
      size=20, color=CINZA)
# 3 side-by-side boxes
for i, (nome, desc, cor) in enumerate([
    ("DOID\n(Diseases)", "12,000+ diseases\norganized by family\n(e.g., diabetes → metabolic disease)", AZUL),
    ("HPO\n(Symptoms)", "16,000+ phenotypes\n(e.g., memory loss,\nseizures, jaundice)", VERDE),
    ("HPOA\n(Associations)", "Which diseases cause\nwhich symptoms —\n320,000+ connections", LARANJA),
]):
    caixa_destaque(sl, f"{nome}\n\n{desc}", 0.4 + i * 4.3, 2.4, 4.1, 3.0,
                   bg=cor, fg=BRANCO, size=17)
texto(sl, "Problem: querying this data requires learning a technical language — SPARQL",
      0.5, 5.6, 12.3, 0.5, size=18, color=VERMELHO, bold=True, align=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 4 — THE BARRIER: SPARQL
# ============================================================
sl = slide_novo()
titulo_slide(sl, "The Technical Barrier: SPARQL")
texto(sl, "To query these databases, a specialist would need to write:", 0.5, 1.75, 12.3, 0.45, size=19, color=CINZA)
caixa_destaque(sl,
    "SELECT ?doenca ?label WHERE {\n"
    "  GRAPH <urn:hpoa> { ?a hpoa:has_phenotype obo:HP_0002354 }\n"
    "  GRAPH <urn:doid> { ?d rdfs:label ?label }\n"
    "  BIND(REPLACE(?xref,\"MIM:\",\"OMIM:\") AS ?omim)\n"
    "}",
    0.5, 2.3, 12.3, 2.0, bg=RGBColor(0x37, 0x47, 0x4F), fg=RGBColor(0xA5, 0xD6, 0xA7),
    size=16, bold=False)
texto(sl, "vs. what the user wants to type:", 0.5, 4.45, 12.3, 0.4, size=19, color=CINZA)
caixa_destaque(sl, '"Which diseases cause memory loss?"',
               0.5, 4.9, 12.3, 0.75, bg=RGBColor(0xE8, 0xF5, 0xE9), fg=VERDE, size=22, bold=False)
texto(sl, "BioSPARQL-NL makes this translation automatically.",
      0.5, 5.8, 12.3, 0.45, size=20, bold=True, color=AZUL, align=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 5 — OUR PROPOSAL
# ============================================================
sl = slide_novo()
titulo_slide(sl, "Our Proposal: BioSPARQL-NL")
# Simple visual flow
for i, (label, cor) in enumerate([
    ("Question\nin natural language", RGBColor(0x42, 0xA5, 0xF5)),
    ("Pipeline\nBioSPARQL-NL", AZUL),
    ("Valid\nSPARQL Query", VERDE),
    ("Database\nAnswer", RGBColor(0x66, 0xBB, 0x6A)),
]):
    caixa_destaque(sl, label, 0.5 + i * 3.2, 2.2, 2.8, 1.4, bg=cor, size=20)
    if i < 3:
        texto(sl, "→", 3.3 + i * 3.2, 2.7, 0.4, 0.6, size=32, bold=True, color=AZUL, align=PP_ALIGN.CENTER)
bullets(sl, [
    "Runs locally on your computer — your data never leaves your machine",
    "Does not require the user to know SPARQL",
    "Automatically validates whether the generated query is correct",
    "Corrects errors and retries if needed",
], y=3.8, size=19, spacing=0.6)

# ============================================================
# SLIDE 6 — ARCHITECTURE
# ============================================================
sl = slide_novo()
titulo_slide(sl, "How It Works: Pipeline Architecture")
img_path = os.path.join(os.path.dirname(__file__), "..", "relatorio", "arquitetura.png")
if os.path.exists(img_path):
    sl.shapes.add_picture(img_path, Inches(0.4), Inches(1.7), Inches(6.5), Inches(5.2))
bullets(sl, [
    "① Extracts entities from the question (NER)",
    "② Retrieves similar examples (few-shot)",
    "③ Builds prompt with context",
    "④ LLM generates the SPARQL query",
    "⑤ Validates against the schema",
    "⑥ Executes in the database",
], x=7.1, y=1.8, w=5.8, size=18, spacing=0.72)

# ============================================================
# SLIDE 7 — NER
# ============================================================
sl = slide_novo()
titulo_slide(sl, "Step 1: Understanding the Question")
texto(sl, "The system uses Artificial Intelligence (scispaCy) to identify\nbiomedical entities in the question:", 0.5, 1.75, 12.3, 0.7, size=19, color=CINZA)
caixa_destaque(sl, '"Which phenotypes are associated with Marfan syndrome?"',
               0.5, 2.6, 12.3, 0.75, bg=RGBColor(0xE3, 0xF2, 0xFD), fg=AZUL, size=21, bold=False)
texto(sl, "↓", 6.3, 3.45, 1.0, 0.5, size=28, bold=True, color=AZUL, align=PP_ALIGN.CENTER)
caixa_destaque(sl, "Recognized: \"Marfan syndrome\"\n→ URI: obo:DOID_0060618  (in the disease database)",
               1.5, 4.05, 10.3, 1.0, bg=RGBColor(0xE8, 0xF5, 0xE9), fg=VERDE, size=19, bold=False)
texto(sl, "This prevents the AI from needing to guess how the disease is registered in the system.",
      0.5, 5.25, 12.3, 0.5, size=17, color=CINZA, italic=True)

# ============================================================
# SLIDE 8 — FEW-SHOT / RAG
# ============================================================
sl = slide_novo()
titulo_slide(sl, "Step 2: Learning from Examples")
texto(sl, "The system maintains a library of 52 questions + already validated queries.", 0.5, 1.75, 12.3, 0.5, size=19, color=CINZA)
texto(sl, "When a new question arrives, it retrieves the 3 most similar ones:", 0.5, 2.3, 12.3, 0.45, size=19, color=CINZA)
for i, (q, sim) in enumerate([
    ('"Symptoms of Alzheimer\'s?"', "92% similar"),
    ('"Phenotypes of Huntington\'s disease?"', "87% similar"),
    ('"Clinical manifestations of multiple sclerosis?"', "81% similar"),
]):
    caixa_destaque(sl, f"{q}  [{sim}]", 0.5, 2.95 + i * 0.82, 12.3, 0.7,
                   bg=RGBColor(0xE8, 0xEF, 0xF8), fg=AZUL, size=17, bold=False)
texto(sl, "These examples are included in the AI prompt — it learns the pattern\nbefore generating the new query. This reduces hallucinations.", 0.5, 5.55, 12.3, 0.7, size=18, color=CINZA)

# ============================================================
# SLIDE 9 — GENERATION
# ============================================================
sl = slide_novo()
titulo_slide(sl, "Step 3: The AI Generates the Query")
texto(sl, "The prompt sent to the AI contains:", 0.5, 1.75, 12.3, 0.45, size=19, color=CINZA)
bullets(sl, [
    "Instructions about the database structure",
    "Schema: which properties and classes exist",
    "Examples of similar questions already solved",
    "The user's question",
], y=2.35, size=18, spacing=0.55)
texto(sl, "Tested models (run locally on your PC):", 0.5, 4.65, 12.3, 0.45, size=19, color=CINZA, bold=True)
for i, (nome, params, cor) in enumerate([
    ("Gemma 3", "4B params", RGBColor(0x42, 0xA5, 0xF5)),
    ("Nemotron Nano", "4B params", VERDE),
    ("Qwen 3.5", "9B params", LARANJA),
    ("Claude Sonnet*", "commercial", VERMELHO),
]):
    caixa_destaque(sl, f"{nome}\n{params}", 0.3 + i * 3.2, 5.2, 2.9, 1.6, bg=cor, size=17)
texto(sl, "* Claude Sonnet ran via API as a performance upper bound reference",
      0.3, 6.9, 12.7, 0.4, size=13, color=CINZA, italic=True)

# ============================================================
# SLIDE 10 — VALIDATION
# ============================================================
sl = slide_novo()
titulo_slide(sl, "Step 4: Validation and Auto-Correction")
texto(sl, "Problem: LLMs sometimes hallucinate properties that do not exist in the database.", 0.5, 1.75, 12.3, 0.5, size=19, color=VERMELHO)
# Validation flow
caixa_destaque(sl, "AI-generated\nSPARQL", 0.3, 2.4, 2.8, 1.2, bg=AZUL, size=17)
texto(sl, "→", 3.1, 2.7, 0.6, 0.6, size=28, bold=True, color=CINZA, align=PP_ALIGN.CENTER)
caixa_destaque(sl, "Schema\nValidation", 3.7, 2.4, 2.8, 1.2, bg=RGBColor(0x78, 0x90, 0x9C), size=17)
texto(sl, "→", 6.5, 2.7, 0.6, 0.6, size=28, bold=True, color=VERDE, align=PP_ALIGN.CENTER)
caixa_destaque(sl, "✓ Valid\n→ execute", 7.1, 2.4, 2.4, 1.2, bg=VERDE, size=17)
texto(sl, "↺", 6.5, 3.7, 0.6, 0.6, size=28, bold=True, color=VERMELHO, align=PP_ALIGN.CENTER)
caixa_destaque(sl, "✗ Invalid\n→ fix (up to 2x)", 7.1, 3.65, 2.4, 1.2, bg=VERMELHO, size=16)
texto(sl, "Validation uses introspection: the system automatically extracts from the database\nwhich properties and classes exist — no manual configuration required.", 0.5, 5.1, 12.3, 0.8, size=18, color=CINZA)
caixa_destaque(sl, "Qwen 9B: from 0% → 50% accuracy with this mechanism alone (+50pp)!",
               0.5, 6.0, 12.3, 0.75, bg=RGBColor(0xE8, 0xF5, 0xE9), fg=VERDE, size=19)

# ============================================================
# SLIDE 11 — REAL EXAMPLE
# ============================================================
sl = slide_novo()
titulo_slide(sl, "Real Working Example")
caixa_destaque(sl, "Question: \"What phenotypes are associated with Marfan syndrome?\"",
               0.3, 1.75, 12.7, 0.7, bg=RGBColor(0xE3, 0xF2, 0xFD), fg=AZUL, size=19, bold=False)
texto(sl, "↓  SPARQL generated and executed:", 0.5, 2.55, 12.3, 0.4, size=17, color=CINZA)
caixa_destaque(sl,
    "SELECT DISTINCT ?phenoLabel WHERE {\n"
    "  GRAPH <urn:hpoa> { ?ann hpoa:source_id \"OMIM:154700\" ;\n"
    "                         hpoa:has_phenotype ?pheno }\n"
    "  GRAPH <urn:hpo> { ?pheno rdfs:label ?phenoLabel }\n"
    "} LIMIT 10",
    0.3, 3.0, 12.7, 1.8, bg=RGBColor(0x37, 0x47, 0x4F), fg=RGBColor(0xA5, 0xD6, 0xA7),
    size=15, bold=False)
caixa_destaque(sl,
    "Results: Arachnodactyly, Ectopia lentis, Aortic root aneurysm,\nScoliosis, Tall stature, Pectus excavatum ...",
    0.3, 4.9, 12.7, 0.9, bg=RGBColor(0xE8, 0xF5, 0xE9), fg=VERDE, size=18, bold=False)
texto(sl, "The user asked in plain English — the system queried 3 databases and returned the phenotypes.",
      0.3, 5.95, 12.7, 0.5, size=17, color=CINZA, italic=True)

# ============================================================
# SLIDE 12 — GOLD STANDARD / EVALUATION
# ============================================================
sl = slide_novo()
titulo_slide(sl, "How We Evaluated the System")
texto(sl, "We created a gold standard: 52 questions with expected answers verified by 2 annotators.", 0.5, 1.75, 12.3, 0.55, size=19, color=CINZA)
for i, (nivel, n, desc, cor) in enumerate([
    ("Easy",   "13 questions", "1 graph, 1 entity\nEx: symptoms of diabetes", RGBColor(0x66, 0xBB, 0x6A)),
    ("Medium", "18 questions", "2 graphs, multiple entities\nEx: diseases with 5+ symptoms", LARANJA),
    ("Hard",   "21 questions", "3 graphs, complex joins,\nspecific MIM/OMIM filters", VERMELHO),
]):
    caixa_destaque(sl, f"{nivel}\n{n}\n\n{desc}", 0.3 + i * 4.35, 2.5, 4.05, 3.2,
                   bg=cor, fg=BRANCO, size=16)
texto(sl, "Main metric: 'correctness' — did the query return the expected results?",
      0.3, 5.9, 12.7, 0.5, size=18, bold=True, color=AZUL, align=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 13 — MODELS
# ============================================================
sl = slide_novo()
titulo_slide(sl, "The 4 AI Models Tested")
tabela_simples(sl,
    ["Model", "Parameters", "Runs on", "Cost"],
    [
        ["Gemma 3 4B",      "4 billion",   "Your PC (8GB GPU)", "Free"],
        ["Nemotron Nano 4B","4 billion",   "Your PC (8GB GPU)", "Free"],
        ["Qwen 3.5 9B",     "9 billion",   "Your PC (CPU+GPU)", "Free"],
        ["Claude Sonnet",   "undisclosed", "Cloud (API)",        "Paid"],
    ],
    x=0.5, y=1.85, col_widths=[3.2, 2.5, 3.3, 2.8],
    row_h=0.6
)
texto(sl, "The first three run locally — no data leaves your computer.\nClaude Sonnet serves as the performance ceiling reference (best possible).", 0.5, 5.45, 12.3, 0.8, size=18, color=CINZA)
caixa_destaque(sl, "Finding: a 4B parameter model outperformed the 9B one — size isn't everything!",
               0.5, 6.35, 12.3, 0.75, bg=RGBColor(0xE8, 0xF5, 0xE9), fg=VERDE, size=18)

# ============================================================
# SLIDE 14 — GENERAL RESULTS 30Q
# ============================================================
sl = slide_novo()
titulo_slide(sl, "Results: 30 Primary Questions")
texto(sl, "With validation enabled (best scenario for each model):", 0.5, 1.75, 12.3, 0.45, size=19, color=CINZA)
# horizontal bars
for i, (label, val, cor) in enumerate([
    ("Nemotron Nano 4B", 80, VERDE),
    ("Qwen 3.5 9B",      50, LARANJA),
    ("Gemma 3 4B",       47, RGBColor(0x42, 0xA5, 0xF5)),
]):
    barra_horizontal(sl, label, val, 100, 0.5, 2.4 + i * 0.95, 9.5, 0.72, cor, size=18)
caixa_destaque(sl, "Surprising finding:\nNemotron 4B (80%) > Qwen 9B (50%)\nArchitecture matters more than size!",
               0.5, 5.35, 12.3, 1.5, bg=AZUL, size=20)

# ============================================================
# SLIDE 15 — VALIDATION IMPACT
# ============================================================
sl = slide_novo()
titulo_slide(sl, "Does Validation Make a Difference?")
texto(sl, "Comparison: with and without the schema-based auto-correction mechanism", 0.5, 1.75, 12.3, 0.45, size=19, color=CINZA)
tabela_simples(sl,
    ["Model", "Without validation", "With validation", "Gain"],
    [
        ["Gemma 3 4B",      "46.7%",  "46.7%",  "0.0pp"],
        ["Nemotron Nano 4B","86.7%",  "80.0%",  "-6.7pp"],
        ["Qwen 3.5 9B",     " 0.0%",  "50.0%",  "+50.0pp"],
    ],
    x=0.5, y=2.3, col_widths=[3.5, 2.8, 2.8, 2.2], row_h=0.65
)
bullets(sl, [
    "Qwen: validation is essential — from 0% to 50% accuracy (without it, useless!)",
    "Gemma: no difference — already generates queries that pass validation",
    "Nemotron: slight reduction — validates correctly but rejects some valid queries",
], y=5.1, size=18, spacing=0.58)

# ============================================================
# SLIDE 16 — BY DIFFICULTY
# ============================================================
sl = slide_novo()
titulo_slide(sl, "Results by Difficulty Level")
tabela_simples(sl,
    ["Model", "Easy (10q)", "Medium (10q)", "Hard (10q)", "Overall"],
    [
        ["Gemma 3 4B",      "60%", "40%", "40%", "47%"],
        ["Nemotron Nano 4B","80%", "80%", "80%", "80%"],
        ["Qwen 3.5 9B",     "50%", "30%", "70%", "50%"],
    ],
    x=0.4, y=1.85, col_widths=[3.2, 2.3, 2.3, 2.3, 2.0], row_h=0.65
)
texto(sl, "Observations:", 0.5, 4.65, 12.3, 0.4, size=18, bold=True, color=AZUL)
bullets(sl, [
    "Nemotron: uniform performance (80% across all levels) — most reliable",
    "Qwen: best on hard questions (70%) — good on complex queries with OPTIONAL/FILTER",
    "Gemma: fastest (4s/query) but less accurate",
], y=5.15, size=17, spacing=0.58)

# ============================================================
# SLIDE 17 — LOCAL VS COMMERCIAL
# ============================================================
sl = slide_novo()
titulo_slide(sl, "Local vs. Commercial: 52 Questions")
texto(sl, "We evaluated Claude Sonnet (commercial reference model) on the full gold standard,\nalongside the best local model (Gemma 3 4B) for comparison:", 0.5, 1.75, 12.3, 0.7, size=18, color=CINZA)
tabela_simples(sl,
    ["Model", "Overall", "Easy", "Medium", "Hard", "Syntax"],
    [
        ["Gemma 3 4B (local)",    "46.2%", "39%", "56%", "43%", "75%"],
        ["Claude Sonnet (cloud)", "82.7%", "100%", "89%", "67%", "100%"],
    ],
    x=0.3, y=2.6, col_widths=[3.5, 1.8, 1.7, 1.7, 1.7, 1.6], row_h=0.65
)
texto(sl, "Claude: 43/52 correct  ·  Gemma: 24/52 correct", 0.3, 4.2, 12.7, 0.45,
      size=17, color=CINZA, align=PP_ALIGN.CENTER)
caixa_destaque(sl, "Gap: 36.5 percentage points", 2.0, 4.75, 9.3, 0.8,
               bg=AZUL, fg=BRANCO, size=22)
texto(sl, "Claude Sonnet has 100% valid syntax — never generates broken SPARQL.",
      0.3, 5.7, 12.7, 0.45, size=17, color=CINZA, italic=True)

# ============================================================
# SLIDE 18 — GAP TRADE-OFF
# ============================================================
sl = slide_novo()
titulo_slide(sl, "The Trade-off: Privacy vs. Precision")
# Gap visualization
retangulo(sl, 0.5, 1.8, 12.3, 1.0, RGBColor(0xE3, 0xF2, 0xFD))
retangulo(sl, 0.5, 1.8, 12.3 * 0.462, 1.0, RGBColor(0x42, 0xA5, 0xF5))
texto(sl, "Gemma local  46.2%", 0.7, 2.0, 5.0, 0.6, size=18, bold=True, color=BRANCO)
retangulo(sl, 0.5, 3.1, 12.3, 1.0, RGBColor(0xE8, 0xF5, 0xE9))
retangulo(sl, 0.5, 3.1, 12.3 * 0.827, 1.0, VERDE)
texto(sl, "Claude Sonnet  82.7%", 0.7, 3.3, 6.0, 0.6, size=18, bold=True, color=BRANCO)
bullets(sl, [
    "Gemma local: data stays on your PC, full privacy — but 36pp less accurate",
    "Claude Sonnet: maximum precision — but your data goes to the cloud (paid API)",
    "For sensitive medical data, the local model may be the right choice",
    "A PC with an 8GB GPU is sufficient to run Gemma 4B",
], y=4.3, size=18, spacing=0.6)

# ============================================================
# SLIDE 19 — ABLATION
# ============================================================
sl = slide_novo()
titulo_slide(sl, "What Does Each Component Contribute?")
texto(sl, "Ablation study: we removed one component at a time and measured impact\n(Gemma 3 4B, 30 questions)", 0.5, 1.75, 12.3, 0.65, size=18, color=CINZA)
tabela_simples(sl,
    ["Configuration", "Syntax", "Accuracy", "Δ vs complete"],
    [
        ["Full pipeline",       "90%", "50.0%", "---"],
        ["Without NER",         "90%", "63.3%", "+13.3pp"],
        ["Without few-shot",    "90%", "46.7%", "-3.3pp"],
        ["Without schema",      "77%", "46.7%", "-3.3pp"],
        ["Without validation",  "67%", "50.0%", "0.0pp"],
        ["Zero-shot (none)",    "87%", "66.7%", "+16.7pp"],
    ],
    x=0.3, y=2.55, col_widths=[3.8, 2.0, 2.0, 3.2], row_h=0.52
)
texto(sl, "* NER (entity recognition) hurts Gemma — injects noise for this specific model.",
      0.3, 6.6, 12.7, 0.45, size=13, color=CINZA, italic=True)

# ============================================================
# SLIDE 20 — FINDING 1: VALIDATION IS CRITICAL
# ============================================================
sl = slide_novo()
titulo_slide(sl, "Finding 1: Validation Is the Critical Component")
texto(sl, "The schema-based auto-correction loop is essential — see the impact without it:", 0.5, 1.75, 12.3, 0.5, size=19, color=CINZA)
for i, (modelo, com_v, sem_v, delta, cor) in enumerate([
    ("Gemma 3 4B",   "46.2%", "25.0%", "−21.2pp", RGBColor(0x42, 0xA5, 0xF5)),
    ("Claude Sonnet","82.7%", "46.2%", "−36.5pp", VERDE),
]):
    retangulo(sl, 0.4, 2.45 + i * 2.2, 12.5, 1.8, RGBColor(0xF5, 0xF7, 0xFA))
    texto(sl, modelo, 0.6, 2.5 + i * 2.2, 3.5, 0.5, size=20, bold=True, color=AZUL)
    caixa_destaque(sl, f"With validation\n{com_v}", 0.5, 2.95 + i * 2.2, 3.8, 1.1, bg=cor, size=19)
    texto(sl, "→", 4.4, 3.25 + i * 2.2, 0.8, 0.6, size=30, bold=True, color=VERMELHO, align=PP_ALIGN.CENTER)
    caixa_destaque(sl, f"Without validation\n{sem_v}", 5.2, 2.95 + i * 2.2, 3.8, 1.1, bg=VERMELHO, size=19)
    caixa_destaque(sl, delta, 9.2, 2.95 + i * 2.2, 3.2, 1.1, bg=RGBColor(0x37, 0x47, 0x4F), size=22)

# ============================================================
# SLIDE 21 — FINDING 2: ARCHITECTURE > SIZE
# ============================================================
sl = slide_novo()
titulo_slide(sl, "Finding 2: Architecture Matters More Than Size")
texto(sl, "Comparing models of different sizes:", 0.5, 1.75, 12.3, 0.45, size=19, color=CINZA)
for i, (nome, params, acerto, cor) in enumerate([
    ("Nemotron Nano", "4 billion", 80, VERDE),
    ("Qwen 3.5",      "9 billion", 50, LARANJA),
    ("Gemma 3",       "4 billion", 47, RGBColor(0x42, 0xA5, 0xF5)),
]):
    barra_horizontal(sl, f"{nome} ({params})", acerto, 100, 0.5, 2.4 + i * 1.05, 9.3, 0.78, cor, size=18)
texto(sl, "Nemotron Nano has 4B parameters (FEWER than Qwen 9B) but achieves 80% vs 50%",
      0.5, 5.7, 12.3, 0.5, size=19, color=CINZA)
caixa_destaque(sl,
    "Conclusion: model architecture and training are more determinant\nthan parameter count for generating correct biomedical SPARQL.",
    0.5, 6.3, 12.3, 0.9, bg=AZUL, size=18)

# ============================================================
# SLIDE 22 — XAI INTERFACE
# ============================================================
sl = slide_novo()
titulo_slide(sl, "Explainable Interface (XAI)")
bullets(sl, [
    "Web interface developed in React + FastAPI",
    "User types the question in natural language",
    "Each pipeline step is shown in real time:",
    "    → Recognized entities (NER)",
    "    → Retrieved examples (few-shot)",
    "    → Generated SPARQL query",
    "    → Validation result and attempts",
    "    → Final database result",
    "Aligned with Explainable AI (XAI) principles: nothing is a 'black box'",
], y=1.85, size=17, spacing=0.5)

# ============================================================
# SLIDE 23 — LIMITATIONS
# ============================================================
sl = slide_novo()
titulo_slide(sl, "What Still Doesn't Work Well")
for i, (titulo, desc, cor) in enumerate([
    ("Hardware as a Bottleneck",
     "Models ≤4B run well on an 8GB GPU. Qwen 9B is 15× slower (63s per query). Models >20B don't fit on common hardware.",
     AZUL),
    ("Idiosyncratic Conventions",
     "9 questions (17%) fail even with Claude Sonnet. They involve technical incompatibilities between databases (e.g., DOID uses 'MIM:' but HPOA uses 'OMIM:').",
     LARANJA),
    ("Stochastic Variability",
     "Gemma 3 4B varied ±3 correct answers across identical runs (temperature=0.1). Results are not 100% deterministic.",
     VERMELHO),
]):
    retangulo(sl, 0.3, 1.85 + i * 1.8, 12.7, 1.6, RGBColor(0xF5, 0xF7, 0xFA))
    caixa_destaque(sl, titulo, 0.3, 1.85 + i * 1.8, 3.2, 1.6, bg=cor, size=16)
    texto(sl, desc, 3.65, 2.0 + i * 1.8, 9.2, 1.3, size=16, color=CINZA)

# ============================================================
# SLIDE 24 — EXPANDED TABLE: 7 MODELS
# ============================================================
sl = slide_novo()
titulo_slide(sl, "Comparison: 7 Models (52 Questions)")
modelos = [
    ("BioGPT 1.5B",      " 0.0%", " 0.0%", "GPT-2 style — no instruction-following"),
    ("BioMedLM 2.7B",    " 0.0%", " 0.0%", "GPT-2 style — no instruction-following"),
    ("Gemma 3 4B",       "46.2%", "75.0%", "General model (baseline)"),
    ("OpenBioLLM-8B",    "57.7%", "61.5%", "Biomedical — weak instruction-following"),
    ("Llama3-Med42-8B",  "67.3%", "90.4%", "Biomedical — strong instruction-following"),
    ("Qwen3-VL-8B",      "78.8%", "98.1%", "General multimodal — best local"),
    ("Claude Sonnet",    "82.7%","100.0%", "Commercial upper bound"),
]
cores_linha = [CINZA, CINZA, AZUL, VERDE, VERDE, LARANJA, VERMELHO]
retangulo(sl, 0.3, 1.3, 12.7, 0.45, AZUL)
for col, cab in enumerate(["Model", "Accuracy", "Syntax", "Note"]):
    xs = [0.35, 5.1, 6.7, 8.1][col]
    ws = [4.6, 1.4, 1.3, 4.8][col]
    texto(sl, cab, xs, 1.32, ws, 0.42, size=13, bold=True, color=BRANCO)
for i, (nome, corr, sint, obs) in enumerate(modelos):
    y = 1.78 + i * 0.72
    bg = RGBColor(0xF0, 0xF4, 0xFF) if i % 2 == 0 else RGBColor(0xFF, 0xFF, 0xFF)
    retangulo(sl, 0.3, y, 12.7, 0.68, bg)
    cor = cores_linha[i]
    texto(sl, nome, 0.35, y + 0.05, 4.6, 0.55, size=13, bold=True, color=cor)
    texto(sl, corr, 5.1,  y + 0.05, 1.4, 0.55, size=13, bold=True, color=cor, align=PP_ALIGN.CENTER)
    texto(sl, sint, 6.7,  y + 0.05, 1.3, 0.55, size=13, color=CINZA, align=PP_ALIGN.CENTER)
    texto(sl, obs,  8.1,  y + 0.05, 4.8, 0.55, size=12, color=CINZA)

# ============================================================
# SLIDE 25 — QALD-9+: GENERALIZATION
# ============================================================
sl = slide_novo()
titulo_slide(sl, "Generalization: QALD-9+ (DBpedia)")
retangulo(sl, 0.4, 1.3, 12.5, 1.5, RGBColor(0xE3, 0xF2, 0xFD))
texto(sl, "Question: does the pipeline work in domains beyond biomedical?",
      0.6, 1.4, 12.0, 0.5, size=17, bold=True, color=AZUL)
texto(sl, "We tested Gemma 3 4B on 30 questions from the public QALD-9+ benchmark (DBpedia) — without domain adaptation.",
      0.6, 1.85, 12.0, 0.45, size=14, color=CINZA)
for i, (label, valor, cor, desc) in enumerate([
    ("Execution OK", "80%",   VERDE,   "SPARQL sent to DBpedia and accepted"),
    ("Avg F1",       "0.106", LARANJA, "Few correct results"),
    ("Exact match",  "3.3%",  VERMELHO,"Only 1/30 perfect"),
]):
    x = 0.5 + i * 4.2
    retangulo(sl, x, 2.95, 3.8, 1.6, RGBColor(0xF5, 0xF5, 0xF5))
    texto(sl, valor, x + 0.15, 3.05, 3.5, 0.9, size=48, bold=True, color=cor, align=PP_ALIGN.CENTER)
    texto(sl, label, x + 0.15, 3.9,  3.5, 0.4, size=15, bold=True, color=AZUL, align=PP_ALIGN.CENTER)
    texto(sl, desc,  x + 0.15, 4.3,  3.5, 0.35, size=12, color=CINZA, align=PP_ALIGN.CENTER)
retangulo(sl, 0.4, 5.1, 12.5, 1.7, RGBColor(0xFF, 0xF3, 0xE0))
texto(sl, "Conclusion:", 0.6, 5.2, 3.0, 0.45, size=15, bold=True, color=LARANJA)
texto(sl, "The pipeline generates valid SPARQL in other domains (80% execution), but needs domain-specific examples and schema\n"
          "to succeed semantically. Syntactic transfer works; semantic transfer does not.",
      0.6, 5.6, 12.0, 1.1, size=14, color=CINZA)

# ============================================================
# SLIDE 26 — CONCLUSIONS
# ============================================================
sl = slide_novo()
titulo_slide(sl, "5 Main Findings")
for i, (letra, titulo, desc, cor) in enumerate([
    ("a", "Validation Is Essential",
     "Gemma: 46% → 25% without it; Claude: 83% → 46% — the auto-correction loop is the most critical component",
     AZUL),
    ("b", "Architecture > Size",
     "Nemotron 4B (80%) beats Qwen 9B (50%) — training matters more than parameter count",
     VERDE),
    ("c", "Gap Nearly Closed with the Right Model",
     "Qwen3-VL 78.8% vs Claude 82.7% (gap of 3.9pp) — next-gen local models enable privacy + precision",
     LARANJA),
    ("d", "Irreducible Errors Exist",
     "9/52 questions (17%) fail even with the best model — intrinsic data limitations",
     VERMELHO),
    ("e", "Biomedical Domain Helps, But Instruction-Following Is Required",
     "Llama3-Med42 (67%) beats Gemma 4B (46%); BioGPT/BioMedLM GPT-2 style = 0%",
     RGBColor(0x6A, 0x1B, 0x9A)),
]):
    retangulo(sl, 0.3, 1.75 + i * 1.08, 0.65, 0.9, cor)
    texto(sl, letra, 0.3, 1.8 + i * 1.08, 0.65, 0.9, size=22, bold=True,
          color=BRANCO, align=PP_ALIGN.CENTER)
    texto(sl, titulo, 1.05, 1.78 + i * 1.08, 3.5, 0.45, size=15, bold=True, color=AZUL)
    texto(sl, desc, 1.05, 2.2 + i * 1.08, 12.0, 0.42, size=13, color=CINZA)

# ============================================================
# SLIDE 27 — THANK YOU
# ============================================================
sl = slide_novo()
retangulo(sl, 0, 0, 13.33, 7.5, AZUL)
retangulo(sl, 0, 4.5, 13.33, 3.0, RGBColor(0x15, 0x26, 0x48))
texto(sl, "Thank You!", 0.6, 1.1, 12.0, 1.3,
      size=60, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
texto(sl, "Questions?", 0.6, 2.5, 12.0, 0.9,
      size=34, color=RGBColor(0xAA, 0xC4, 0xFF), align=PP_ALIGN.CENTER)
texto(sl, "Cassiano Ricardo Neubauer Moralles",
      0.6, 4.7, 12.0, 0.5, size=18, color=BRANCO, align=PP_ALIGN.CENTER)
texto(sl, "cass78@gmail.com  |  Graduate Program in Applied Computing — UNISINOS",
      0.6, 5.25, 12.0, 0.45, size=15, color=RGBColor(0xAA, 0xC4, 0xFF), align=PP_ALIGN.CENTER)
texto(sl, "Code and data: output/eval_*.json  |  output/ablation_*.json  |  output/relatorio/biosparql-nl.pdf",
      0.6, 5.75, 12.0, 0.4, size=12, color=RGBColor(0x78, 0x90, 0x9C), align=PP_ALIGN.CENTER)

# ============================================================
out = os.path.join(os.path.dirname(__file__), "biosparql-nl-apresentacao-en.pptx")
prs.save(out)
print(f"[OK] Presentation saved at: {out}")
print(f"     Total slides: {len(prs.slides)}")
